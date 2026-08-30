# SPDX-License-Identifier: FSL-1.1-Apache-2.0
# Copyright (c) 2025 Open Computer Use Contributors
"""GROUP I — MCP tool-surface properties BELOW the exec contract, fleet-only.

The @L4 half of mcp_tool_surface.feature (projection, script semantics, result
shaping) is proven in the gateway forward e2e, where a control-mock executes the
real committed scripts. This group proves the @L5 half — the properties only a
LIVE guest can settle:

  I1  a real failing command's non-zero exit is really transported to isError
  I2  oversized guest output is bounded by control, not relayed whole
  I3  a command past the timeout is killed below the contract
  I4  create_file EACCES is a function of the live guest identity (a contrast)
  I5  one chat maps to one PERSISTENT workspace across tool calls
  I6  the four tools compose over one workspace (needs a python3 guest, #122)

It drives the REAL gateway on 127.0.0.1:8080 with the minted bearer, exactly like
group H. Every scenario carries its keystone (the negative that makes the green
non-vacuous) inline. It skips loudly when the gateway is unreachable or the boot
set is not rendered — it never fabricates a green.

No poc counterpart: these are gateway/control-transported properties; the PoC has
no gateway. The PoC-vs-fleet contrast for each lives in scenarios.yaml (I1..I6).
"""

import json
import os
import subprocess
import time
import uuid
from pathlib import Path

import pytest

# The gateway's north listener, host-mapped in the fleet compose (same as H).
GATEWAY_URL = "http://127.0.0.1:8080/"
_PROTO = "2025-06-18"
# The secrets dir has TWO sources, in priority order:
#   1. the dir the RUNNING gateway container actually mounts its boot-set from
#      (resolved live via docker inspect) -- the bearer there is the ONLY one that
#      hashes to the gateway's loaded credential, so this is authoritative;
#   2. the file-relative fallback (this checkout's own fleet/secrets/gateway).
# Preferring the mounted dir removes the wrong-stage-tree trap: a suite run from a
# sibling checkout still reads the bearer the live stack accepts, instead of a
# stale sibling bearer that 401s (which used to read as a clean capability skip).
_GATEWAY_CONTAINER = os.getenv("OCU_GATEWAY_CONTAINER", "ocu-donegate-mcp-gateway-1")
_SECRETS_FALLBACK = Path(__file__).resolve().parents[2] / "fleet" / "secrets" / "gateway"


def _running_gateway_secrets_dir():
    """The dir the live gateway container mounts boot-set.json from, or None.

    Reads the container's mounts via docker inspect; returns the parent dir of the
    boot-set mount SOURCE so the matching plaintext bearer.txt (minted alongside it)
    is read from the same place the gateway validates against. None if docker is
    unavailable, the container is absent, or no boot-set mount is found.
    """
    try:
        out = subprocess.run(
            ["docker", "inspect", "--format", "{{json .Mounts}}", _GATEWAY_CONTAINER],
            capture_output=True, text=True, timeout=10,
        )
        if out.returncode != 0:
            return None
        for m in json.loads(out.stdout):
            if m.get("Destination", "").endswith("/boot-set.json"):
                src = m.get("Source")
                if src:
                    return Path(src).parent
    except (OSError, ValueError, subprocess.SubprocessError):
        return None
    return None


def _secrets_dir():
    d = _running_gateway_secrets_dir()
    if d is not None and (d / "bearer.txt").exists():
        return d
    return _SECRETS_FALLBACK


_SECRETS = _secrets_dir()
_BOOT_SET = _SECRETS / "boot-set.json"
_BEARER_FILE = _SECRETS / "bearer.txt"

pytestmark = pytest.mark.fleet


def _bearer():
    if _BEARER_FILE.exists():
        b = _BEARER_FILE.read_text().strip()
        return b or None
    return None


def _bash_body(command):
    """A tools/call bash_tool JSON-RPC body — the exact shape the browser sends."""
    return json.dumps(
        {
            "jsonrpc": "2.0",
            "id": 1,
            "method": "tools/call",
            "params": {"name": "bash_tool", "arguments": {"command": command}},
        }
    )


def _file_tool_body(name, arguments):
    """A tools/call body for a file tool (create_file / view / str_replace)."""
    return json.dumps(
        {
            "jsonrpc": "2.0",
            "id": 1,
            "method": "tools/call",
            "params": {"name": name, "arguments": arguments},
        }
    )


def _call(chat_id, body, timeout=40):
    """POST a tools/call to the gateway; return (http_status:int, parsed:dict|None).

    Uses curl to match the exact wire and to stay dependency-free like the demo
    scripts. Raises on a transport failure so an unreachable gateway becomes a
    loud skip via _require_gateway, never a silent pass.
    """
    bearer = _bearer()
    args = [
        "curl", "-sS", "--max-time", str(timeout), "-o", "-", "-w", "\n%{http_code}",
        GATEWAY_URL,
        "-H", f"Authorization: Bearer {bearer}",
        "-H", f"MCP-Protocol-Version: {_PROTO}",
        "-H", f"X-Chat-Id: {chat_id}",
        "-H", "content-type: application/json",
        "-d", body,
    ]
    out = subprocess.run(args, capture_output=True, text=True, timeout=timeout + 10)
    if out.returncode != 0:
        raise RuntimeError(f"curl transport failure rc={out.returncode}: {out.stderr[:200]}")
    text = out.stdout
    nl = text.rfind("\n")
    status = int(text[nl + 1 :].strip())
    try:
        parsed = json.loads(text[:nl])
    except (ValueError, TypeError):
        parsed = None
    return status, parsed


def _result(parsed):
    """(text, is_error) from a CallToolResult, or (None, None) on a JSON-RPC error."""
    if not isinstance(parsed, dict) or "error" in parsed:
        return None, None
    res = parsed.get("result") or {}
    is_error = bool(res.get("isError"))
    text = None
    for item in res.get("content") or []:
        if item.get("type") == "text":
            text = item.get("text")
            break
    return text, is_error


def _gateway_live():
    try:
        _call("i-probe", _bash_body("true"), timeout=8)
        return True
    except Exception:
        return False


def _guest_has_python3(chat_id):
    """True iff the live guest has a runnable python3 (gates the file-tool journey).

    The probe command echoes "yes" or "no", so a non-empty reply is authoritative.
    An EMPTY reply is not "no python3" -- it is a transient exec miss (a session-
    create hiccup late in a long suite returns 200 with no content block, which used
    to skip the skill tests as if the guest were thin). Retry once on empty before
    deciding; a genuine "no" is returned immediately, never retried.
    """
    probe = _bash_body("command -v python3 >/dev/null 2>&1 && echo yes || echo no")
    for _ in range(2):
        _, parsed = _call(chat_id, probe)
        text, _err = _result(parsed)
        if text and text.strip():
            return "yes" in text
        # empty reply -> transient exec miss; retry once with a fresh chat id so a
        # per-session create race does not repeat.
        chat_id = f"{chat_id}-r"
    return False


@pytest.fixture(autouse=True)
def _require_gateway():
    if not _BOOT_SET.exists():
        pytest.skip(
            f"gateway boot-set not rendered at {_BOOT_SET} "
            "(run scripts/mint_boot_set.py) — SKIP, not a pass."
        )
    if _bearer() is None:
        pytest.skip(f"gateway bearer not rendered at {_BEARER_FILE} — SKIP, not a pass.")
    if not _gateway_live():
        pytest.skip(
            f"MCP gateway not reachable at {GATEWAY_URL} "
            "(bring up the mcp-gateway service inside Lima) — SKIP, not a pass."
        )
    # A REACHABLE gateway that rejects the bearer is a harness auth desync, never a
    # skip: the bearer file and the running gateway's boot-set were minted in
    # different runs / different stage trees, so the presented sk-ocu- token does not
    # hash to the gateway's stored credential. Silently skipping here let a stale
    # bearer masquerade as a clean capability skip (the guest probe returns empty on
    # the 401, reading as "no python3"). FAIL loudly so a wrong-tree run is caught.
    status, parsed = _call("i-auth-probe", _bash_body("true"), timeout=8)
    if status == 401 or (isinstance(parsed, dict) and (parsed.get("error") or {}).get("message") == "unauthenticated"):
        pytest.fail(
            f"gateway at {GATEWAY_URL} returned 401 unauthenticated for the bearer at "
            f"{_BEARER_FILE}: the bearer does not match the running gateway's boot-set "
            "(run the suite from the stage tree whose secrets the live stack mounts, or "
            "re-mint the boot-set). A reachable-but-401 gateway is a harness desync, not a skip."
        )


def _cid(tag):
    return f"i-{tag}-{uuid.uuid4().hex[:12]}"


# ---------------------------------------------------------------------------
# I1 — a real non-zero exit is transported to isError (not fabricated)
# ---------------------------------------------------------------------------

def test_i1_nonzero_exit_transports_to_iserror():
    """A real failing command over the live gateway → isError:true; a zero-exit
    command over the SAME path → isError:false. Proves the flag tracks the guest's
    real exit code (control transports it), not a constant. Empty-output non-zero
    also carries the synthesized "[Exit code: N]" (the L4 shaping, proven live)."""
    status, parsed = _call(_cid("i1-fail"), _bash_body("exit 7"))
    assert status == 200, f"a tool error is Tier-2 (HTTP 200 + isError), got {status}"
    text, is_error = _result(parsed)
    assert is_error is True, f"a real non-zero exit must transport isError:true, got {parsed}"
    assert text == "[Exit code: 7]", (
        f"empty-output non-zero must synthesize the exit-code marker live, got {text!r}"
    )

    # keystone: a zero exit over the SAME path is NOT an error (flag is not constant).
    status, parsed = _call(_cid("i1-ok"), _bash_body("true"))
    text, is_error = _result(parsed)
    assert status == 200 and is_error is False, (
        f"a zero-exit command must be isError:false (keystone), got status={status} {parsed}"
    )


# ---------------------------------------------------------------------------
# I2 — oversized output is bounded by control, not relayed whole
# ---------------------------------------------------------------------------

# The caller-facing content ceiling: control bounds each F5 reply stream at 64KiB at
# the source (#128), so the DATA a caller sees is <= 64KiB. A truncation marker is
# appended AFTER that bound, so the total text is 64KiB of data + a short marker.
_CONTENT_CEILING_BYTES = 64 << 10  # 65536
_MARKER_HEADROOM = 128  # the "[output truncated at N bytes]" marker rides on top of the ceiling

_I2_BIG = 120000  # raw stdout for the oversized probe — well past the 64KiB ceiling


def test_i2_oversized_output_bounded_at_ceiling_with_marker():
    """The merged #127+#128 contract (was I2 nothing-lost XOR I2c ceiling, now one
    bounded-contract test after the paired flip): an oversized output is a BOUNDED
    TOOL RESULT — HTTP 200, isError:false, DATA truncated to the 64KiB caller ceiling
    with a truncation marker appended — NOT a 502 that loses the whole result (#127)
    and NOT relayed whole (#128 bounds it at the source).

    Live-verified: a 120k output that once 502'd, then came back whole at 8MiB, now
    returns 64KiB of data + '[output truncated at N bytes]'. The DATA (before the
    marker) is <= 64KiB; the total is that plus a short marker."""
    status, parsed = _call(_cid("i2-big"), _bash_body(f"yes X | head -c {_I2_BIG}"))
    text, is_error = _result(parsed)
    assert status == 200, (
        f"oversized output must be a bounded tool result (HTTP 200), not a 502 "
        f"forward-refusal that loses the whole result — got {status}"
    )
    assert is_error is False, f"a truncated large output is not an error, got isError on {parsed}"
    assert text is not None, f"oversized output must return content, got {parsed}"
    assert "truncat" in text.lower(), (
        f"a bounded large output must carry a truncation marker, got {text[-120:]!r}"
    )
    # The DATA (everything before the appended marker) is bounded at the 64KiB ceiling;
    # the marker rides on top, so the total stays within a small headroom of it.
    assert len(text) <= _CONTENT_CEILING_BYTES + _MARKER_HEADROOM, (
        f"oversized output must be bounded at the 64KiB ceiling (+ a short marker), got "
        f"{len(text)} bytes (ceiling {_CONTENT_CEILING_BYTES} + {_MARKER_HEADROOM} headroom)"
    )
    # ...and it is genuinely bounded BELOW the emitted size (nothing silently un-bounded).
    assert len(text) < _I2_BIG, (
        f"a {_I2_BIG}-byte emit must come back bounded, got {len(text)} bytes"
    )

    # keystone: a small marker over the SAME path comes back whole, un-truncated.
    marker = "I2_SMALL_" + uuid.uuid4().hex[:8]
    _, parsed = _call(_cid("i2-small"), _bash_body(f"printf %s {marker}"))
    text, _ = _result(parsed)
    assert text is not None and marker in text and "truncat" not in text.lower(), (
        f"a small output must return whole and un-truncated (keystone), got {text!r}"
    )


def test_i2b_moderate_output_returns_whole():
    """A moderate output returns whole and un-truncated over the live gateway. This
    is the long-lived whole-return-under-ceiling keystone: n=30000 sits below the
    FINAL 64KiB caller ceiling (#128), so it returns whole both now (step 1: 8MiB
    boundContent) and after step 2 lands (64KiB ceiling) — it survives the step-2
    flip untouched. Kept deliberately below the ceiling; do not raise toward 64KiB."""
    n = 30000
    status, parsed = _call(_cid("i2b"), _bash_body(f"yes X | head -c {n}"))
    text, is_error = _result(parsed)
    assert status == 200 and is_error is False, (
        f"a moderate {n}-byte output must return a clean tool result, got {status} {parsed}"
    )
    assert text is not None and len(text) >= n, (
        f"a moderate output must come back whole ({n} bytes), got {len(text or '')}"
    )


# ---------------------------------------------------------------------------
# I3 — a command past the timeout is killed below the contract
# ---------------------------------------------------------------------------

def test_i3_timeout_is_enforced():
    """A command that would sleep far past the exec-timeout is KILLED (does not run
    to completion), while a fast command over the SAME path returns promptly. This
    proves timeout ENFORCEMENT — the load-bearing property — holds on the live
    stand. It does NOT assert the result shape (that is I3b).

    The contract: the gateway provisioning policy sets exec_timeout_seconds=120
    (deploy/fleet/secrets/gateway/provisioning-policy.json, the PoC-parity value),
    bounded by control's host-side hard cap of 5 minutes. The kill therefore lands
    at ~120s; the transport budget must OUTLIVE the kill (a curl budget equal to
    the kill time aborts rc=28 with zero bytes and reads as a false red)."""
    # sleep 600 would hang for 10 minutes if unbounded; the exec-timeout kills it
    # at ~120s. We assert it did NOT run to completion, whatever the result shape.
    start = time.monotonic()
    try:
        _call(_cid("i3-hang"), _bash_body("sleep 600"), timeout=180)
    except RuntimeError:
        # A curl-level abort still means the call did not hang the full 600s; the
        # wall-clock assertion below is the real check.
        pass
    elapsed = time.monotonic() - start
    assert elapsed < 150, (
        f"a sleep 600 must be KILLED by the 120s exec-timeout (policy "
        f"exec_timeout_seconds), not run to completion; took {elapsed:.0f}s "
        f"(timeout enforcement is the load-bearing property)"
    )

    # keystone: a fast command returns promptly (the kill is the timeout, not a cap).
    start = time.monotonic()
    marker = "I3_FAST_" + uuid.uuid4().hex[:8]
    _, parsed = _call(_cid("i3-fast"), _bash_body(f"printf %s {marker}"))
    fast_elapsed = time.monotonic() - start
    text, _ = _result(parsed)
    assert text is not None and marker in text and fast_elapsed < 30, (
        f"a fast command must return promptly (keystone), took {fast_elapsed:.0f}s text={text!r}"
    )


def test_i3b_timeout_surfaces_as_tool_result_with_partial_output():
    """#129 (FIXED, control PR #63): a timed-out command is a Tier-2 tool result
    (HTTP 200 + isError) carrying its PARTIAL output and a timeout notice, the way
    the PoC returns it — not a 502 that loses the whole result. Companion to I3: I3
    proves the KILL; I3b pins the RESULT SHAPE (exit-124 reply, partial output
    preserved, notice appended).

    Live-verified: control shapes a host exec-timeout into a valid exit-124 reply
    with the pre-kill stdout + '[Command timed out after Ns]' in the SAME stream the
    gateway relays on isError (stdout AND non-empty stderr), so the partial output
    survives the gateway's stderr-wins relay."""
    # A command that writes a marker BEFORE it hangs past the exec-timeout. The
    # transport budget must outlive the ~120s policy kill plus the result relay,
    # or curl aborts rc=28 at exactly the kill boundary and loses the reply.
    marker = "I3B_PARTIAL_" + uuid.uuid4().hex[:8]
    status, parsed = _call(_cid("i3b"), _bash_body(f"echo {marker}; sleep 600"), timeout=180)
    text, is_error = _result(parsed)
    assert status == 200, (
        f"a timed-out command must be a Tier-2 tool result (HTTP 200 + isError), not "
        f"a 502 forward-refusal that loses the result — got {status}"
    )
    assert is_error, f"a timed-out command must be a tool error (exit 124), got {parsed}"
    assert text is not None and marker in text, (
        f"the PARTIAL output written before the kill must survive to the caller "
        f"(#129: it used to be lost to the 502), got {text!r}"
    )
    assert "timed out" in text.lower(), (
        f"a timed-out result must carry the timeout notice, got {text!r}"
    )


# ---------------------------------------------------------------------------
# I4 — create_file EACCES is a function of the live guest identity (a contrast)
# ---------------------------------------------------------------------------

def test_i4_write_permission_is_guest_identity_contrast():
    """The PoC's create_file into /home/assistant returns [Errno 13] because the
    guest user is not the owner. The fleet guest identity may differ (root vs a
    named user), so this records the LIVE outcome as a contrast rather than
    asserting the PoC's exact errno. The keystone: a writable path over the SAME
    session succeeds — the write plane is not broken, only the protected path is
    (whatever the guest identity makes protected)."""
    if not _guest_has_python3(_cid("i4-probe")):
        pytest.skip(
            "create_file needs a python3-bearing guest (#122); the default guest is "
            "stripped — SKIP, not a pass. I4 is a file-tool identity contrast."
        )
    cid = _cid("i4")
    # A path a non-root user typically cannot write; the LIVE outcome is recorded.
    status, parsed = _call(
        cid, _file_tool_body("create_file", {"path": "/root/i4_probe.txt", "file_text": "x"})
    )
    text, is_error = _result(parsed)
    assert status == 200, f"create_file is a tool call (HTTP 200), got {status}"
    # Contrast, not a fixed assertion: if the guest runs as root the write may
    # SUCCEED where the PoC's assistant-user would EACCES. Either way it must be a
    # legible result, and we record which for CONTRAST.md.
    outcome = "ERROR(" + (text or "").strip()[:60] + ")" if is_error else "WROTE"
    print(f"[I4 CONTRAST] fleet create_file /root -> {outcome}")

    # keystone: a writable path over the SAME session must succeed (write plane OK).
    marker = "I4_OK_" + uuid.uuid4().hex[:8]
    _, parsed = _call(
        cid, _file_tool_body("create_file", {"path": f"/tmp/{marker}.txt", "file_text": marker})
    )
    text, is_error = _result(parsed)
    assert is_error is False and text and "Successfully created" in text, (
        f"a writable /tmp path must succeed (keystone — write plane not broken), got {text!r}"
    )


# ---------------------------------------------------------------------------
# I5 — one chat maps to one PERSISTENT workspace across tool calls
# ---------------------------------------------------------------------------

def test_i5_workspace_persists_across_calls_in_one_session():
    """Call 1 writes a marker file; call 2 in the SAME chat reads it back. Proves
    the session maps to ONE persistent container, not a fresh one per exec. The
    keystone: a DIFFERENT chat id does NOT see the marker — persistence is
    per-session, not a shared global filesystem. This is sh-only (no python3),
    provable on the default guest — the highest-value single property here."""
    cid = _cid("i5")
    marker = "I5_PERSIST_" + uuid.uuid4().hex[:12]

    # call 1 — write the marker.
    _, parsed = _call(cid, _bash_body(f"echo {marker} > /tmp/i5_probe.txt && echo wrote"))
    text, is_error = _result(parsed)
    assert is_error is False and text and "wrote" in text, (
        f"call 1 must write the marker file, got {text!r}"
    )

    # call 2 — same chat reads it back.
    _, parsed = _call(cid, _bash_body("cat /tmp/i5_probe.txt"))
    text, is_error = _result(parsed)
    assert is_error is False and text and marker in text, (
        f"call 2 in the SAME session must see the marker call 1 wrote (workspace persisted), "
        f"got {text!r}"
    )

    # keystone: a DIFFERENT chat must NOT see the marker (per-session isolation).
    other = _cid("i5-other")
    _, parsed = _call(other, _bash_body("cat /tmp/i5_probe.txt 2>&1 || true"))
    text, _ = _result(parsed)
    assert text is None or marker not in text, (
        f"a DIFFERENT chat must NOT see another session's file (keystone — per-session, "
        f"not shared global fs), got {text!r}"
    )


# ---------------------------------------------------------------------------
# I6 — the four tools compose over one workspace (the owner-shown journey)
# ---------------------------------------------------------------------------

def test_i6_four_tools_compose_over_one_workspace():
    """create_file writes a file, view shows it, str_replace edits it, bash_tool
    cats it — all over ONE session. The full owner-shown journey. Requires a
    python3-bearing guest (#122); skips loudly on a stripped guest rather than
    passing vacuously. No tool may fall back to bash for a file operation."""
    probe_cid = _cid("i6-probe")
    if not _guest_has_python3(probe_cid):
        pytest.skip(
            "the four-tool journey needs a python3-bearing guest (#122); the default "
            "guest is stripped — SKIP, not a pass."
        )
    cid = _cid("i6")
    path = "/tmp/i6_journey.txt"

    # 1) create_file — ALPHA original
    _, parsed = _call(cid, _file_tool_body("create_file", {"path": path, "file_text": "ALPHA original\n"}))
    text, is_error = _result(parsed)
    assert is_error is False and text and "Successfully created" in text, (
        f"create_file must write the file, got {text!r}"
    )

    # 2) view — shows the file with line numbers
    _, parsed = _call(cid, _file_tool_body("view", {"path": path}))
    text, is_error = _result(parsed)
    assert is_error is False and text and "ALPHA original" in text, (
        f"view must show the created file, got {text!r}"
    )

    # 3) str_replace — original -> EDITED
    _, parsed = _call(
        cid, _file_tool_body("str_replace", {"path": path, "old_str": "original", "new_str": "EDITED"})
    )
    text, is_error = _result(parsed)
    assert is_error is False and text and "Successfully replaced" in text, (
        f"str_replace must edit the file, got {text!r}"
    )

    # 4) bash_tool — cat confirms the edited content over the SAME workspace
    _, parsed = _call(cid, _bash_body(f"cat {path}"))
    text, is_error = _result(parsed)
    assert is_error is False and text and "ALPHA EDITED" in text, (
        f"bash cat must confirm the edit persisted across all four tools, got {text!r}"
    )


def _image_url_block(parsed):
    """The image content block (data-URI) from a view result, or None.

    view on an IMAGE must return a real image the MODEL sees (D4, PARITY-LEDGER
    -147): the gateway ViewScript emits a resized JPEG and projectCallToolResult
    wraps it in an image_url content block. This reaches past _result (which only
    reads the text block) into res.content for the image_url item.
    """
    if not isinstance(parsed, dict):
        return None
    for item in (parsed.get("result") or {}).get("content") or []:
        if isinstance(item, dict) and item.get("type") == "image_url":
            url = item.get("image_url")
            return url.get("url") if isinstance(url, dict) else url
    return None


def test_i7_view_on_image_returns_a_model_visible_image_block():
    """D4 parity: view on an IMAGE gives the MODEL a rendered image block (a
    data:image data-URI), not just the numbered-line text path. This is the
    tool-side of image handling, distinct from the pane preview (M1): the model
    itself sees the pixels via the view tool. Requires a PIL-bearing guest.

    Non-vacuous: assert a genuine data:image data-URI content block, not that the
    text merely mentions the filename. A guest with no image projection returns
    only a text block -> the image_url is None -> RED.
    """
    probe_cid = _cid("i7-probe")
    if not _guest_has_python3(probe_cid):
        pytest.skip(
            "view-image needs a PIL-bearing guest (#122); the default guest is "
            "stripped -- SKIP, not a pass."
        )
    cid = _cid("i7")
    path = "/tmp/i7_view.png"

    # 1) the guest writes a small PNG of known dimensions (avoid nested quotes:
    # a heredoc into the guest's python, chr() for the mode string).
    make = (
        "python3 - <<'PY'\n"
        "from PIL import Image\n"
        f"Image.new(chr(82)+chr(71)+chr(66), (64, 48), (7, 8, 9)).save({path!r})\n"
        "print('WROTE')\n"
        "PY"
    )
    _, parsed = _call(cid, _bash_body(make))
    text, is_error = _result(parsed)
    assert is_error is False and text and "WROTE" in text, (
        f"guest PNG write for the view leg failed: {text!r}"
    )

    # 2) view the image -> the result must carry a model-visible image block.
    _, parsed = _call(cid, _file_tool_body("view", {"path": path}))
    _, is_error = _result(parsed)
    assert is_error is False, f"view on an image errored: {parsed!r}"
    url = _image_url_block(parsed)
    assert url and str(url).startswith("data:image"), (
        "view on an image must return a data:image data-URI content block the "
        f"model can see (D4 parity); got image_url={url!r}. A text-only result "
        "means the guest did not project the image."
    )


def test_i8_document_skill_produces_a_valid_ooxml_artifact_in_guest():
    """Parity: the guest's document-skill toolchain (pandoc markdown -> docx)
    PRODUCES a real, content-bearing OOXML artifact -- the non-image, non-inline
    "skills produce their artifact" leg. Distinct from M3 (matplotlib PNG,
    image-previewable)
    and the B-group (which validates the download path with a TEST-BUILT docx).
    Here the ARTIFACT ITSELF is produced by the skill inside the guest.

    Non-vacuous: the guest base64s the produced docx out, the test unzips it and
    asserts BOTH word/document.xml is present (valid OOXML) AND the unique body
    marker is in it (real content, not an empty/garbage zip). A guest with no
    pandoc, or a corrupt artifact, reds. Needs a pandoc-bearing guest.
    """
    import base64
    import io
    import zipfile

    cid = _cid("i8")
    if not _guest_has_python3(cid):
        pytest.skip("document skill needs a full guest (#122) -- SKIP, not a pass.")

    marker = "DOCXMARK" + uuid.uuid4().hex[:8]
    out = "/tmp/i8_report.docx"
    src = f"# Report\nThe unique body text is {marker} here.\n"
    # Build the source, run the skill (pandoc), stream the artifact out as base64.
    gcmd = (
        "printf %s " + repr(src) + " > /tmp/i8_src.md && "
        "pandoc /tmp/i8_src.md -o " + out + " && "
        "base64 -w0 " + out
    )
    _, parsed = _call(cid, _bash_body(gcmd), timeout=150)
    text, is_error = _result(parsed)
    if is_error and text and "pandoc" in text and "not found" in text:
        pytest.skip("guest has no pandoc (document skill absent) -- SKIP, not a pass.")
    assert is_error is False and text, f"document skill run failed: {text!r}"

    b64 = next((ln for ln in reversed(text.splitlines()) if ln.strip()), "")
    assert len(b64) > 100, f"no base64 artifact streamed from the guest: {text!r}"
    data = base64.b64decode(b64)

    z = zipfile.ZipFile(io.BytesIO(data))
    assert "word/document.xml" in z.namelist(), (
        "the skill artifact is not a valid OOXML docx (no word/document.xml)"
    )
    doc_xml = z.read("word/document.xml").decode()
    assert marker in doc_xml, (
        f"the produced docx does not carry the source content ({marker!r} absent "
        "from word/document.xml) -- the skill did not render the real document"
    )


def test_i9_spreadsheet_skill_produces_a_valid_xlsx_with_data_in_guest():
    """Parity: the guest's spreadsheet-skill toolchain (openpyxl) PRODUCES a
    real, data-bearing .xlsx -- a DIFFERENT OOXML type (xl/worksheets/sheet1.xml)
    and a DIFFERENT skill runtime than i8's pandoc docx, exercising the data/
    spreadsheet artifact the model builds. The test reads the artifact back with
    zipfile + raw XML only (no openpyxl needed on the test side).

    Non-vacuous: assert the produced workbook has a real worksheet part AND the
    unique cell value is in it (not an empty/garbage zip). A guest with no
    openpyxl, or a workbook missing the value, reds. Needs a full guest.
    """
    import base64
    import io
    import zipfile

    cid = _cid("i9")
    if not _guest_has_python3(cid):
        pytest.skip("spreadsheet skill needs a full guest (#122) -- SKIP, not a pass.")

    cell = "XLCELL" + uuid.uuid4().hex[:8]
    out = "/tmp/i9_book.xlsx"
    # The skill: openpyxl writes a workbook with a known cell value, streamed out.
    # chr() avoids nested quoting for the A1 cell reference inside the heredoc.
    gcmd = (
        "python3 - <<PY\n"
        "import openpyxl\n"
        "wb = openpyxl.Workbook(); ws = wb.active\n"
        f"ws[chr(65) + chr(49)] = {cell!r}\n"
        f"wb.save({out!r})\n"
        "PY\n"
        "base64 -w0 " + out
    )
    _, parsed = _call(cid, _bash_body(gcmd), timeout=150)
    text, is_error = _result(parsed)
    if is_error and text and "openpyxl" in text and "No module" in text:
        pytest.skip("guest has no openpyxl (spreadsheet skill absent) -- SKIP, not a pass.")
    assert is_error is False and text, f"spreadsheet skill run failed: {text!r}"

    b64 = next((ln for ln in reversed(text.splitlines()) if ln.strip()), "")
    assert len(b64) > 100, f"no base64 artifact streamed from the guest: {text!r}"
    data = base64.b64decode(b64)

    z = zipfile.ZipFile(io.BytesIO(data))
    names = z.namelist()
    assert any("xl/worksheets/sheet1.xml" in n for n in names), (
        f"the skill artifact is not a valid xlsx (no worksheet part); parts={names}"
    )
    sheet_xml = z.read("xl/worksheets/sheet1.xml").decode(errors="ignore")
    assert cell in sheet_xml, (
        f"the produced xlsx does not carry the cell value ({cell!r} absent from "
        "the worksheet) -- the skill did not write the real data"
    )


def test_i10_str_replace_edits_a_file_on_the_outputs_surface():
    """Parity + regression: str_replace edits an EXISTING file on the object-
    backed /mnt/user-data/outputs surface, and the edit PERSISTS (visible to a
    FRESH guest session, so it is not merely a write-back cache artifact of the
    editing session).

    This guards a real defect: the projection scripts used open(path,'w')
    (O_TRUNC-on-open), which the object-backed FUSE outputs surface rejects with
    EIO for an EXISTING file (create_file only wrote NEW files so it never
    tripped it; str_replace edits existing so it always did). The fix is a single
    r+ seek/write/truncate. Non-vacuous by two independent asserts:
      1. the tool result is "Successfully replaced" and NOT an "Error:" line
         (the broken pattern returns "Error: [Errno 5] Input/output error");
      2. a FRESH session reads the edited content back through the surface
         (proves persistence, not a same-session cache echo).
    Needs a python3-bearing guest.
    """
    # ONE chat_id (== one derived per-chat storage scope, D5). The outputs surface
    # is scoped per-chat, so persistence is verified WITHIN the same scope by a
    # later exec -- a different chat_id would be a different scope and correctly
    # would NOT see the file (that is the D5 isolation contract, not persistence).
    editor = _cid("i10")
    if not _guest_has_python3(editor):
        pytest.skip("str_replace-on-outputs needs a full guest (#122) -- SKIP, not a pass.")

    marker_before = "I10_BEFORE_" + uuid.uuid4().hex[:8]
    marker_after = "I10_AFTER_" + uuid.uuid4().hex[:8]
    name = "i10_" + uuid.uuid4().hex[:8] + ".txt"
    path = "/mnt/user-data/outputs/" + name

    # 1) create the file on the outputs surface (exists as an object).
    _, parsed = _call(
        editor, _file_tool_body("create_file", {"path": path, "file_text": marker_before + " original\n"})
    )
    text, is_error = _result(parsed)
    assert is_error is False and text and "Successfully created" in text, (
        f"create_file onto outputs failed: {text!r}"
    )

    # 2) str_replace the EXISTING outputs file -- this is the O_TRUNC-on-existing
    # path that EIO'd before the r+ fix.
    _, parsed = _call(
        editor,
        _file_tool_body("str_replace", {"path": path, "old_str": marker_before, "new_str": marker_after}),
    )
    text, is_error = _result(parsed)
    assert is_error is False and text and "Successfully replaced" in text and "Error" not in text, (
        f"str_replace on the outputs surface failed: {text!r}. The broken "
        "open(path,'w') pattern returns '[Errno 5] Input/output error' here; the "
        "r+ seek/write/truncate fix restores it."
    )

    # 3) a LATER exec in the SAME scope reads the edited content back through the
    # surface -- proves the edit PERSISTED to the object, not a lost/torn write.
    # A separate guest exec (cat) forces a fresh open of the object, so a torn or
    # zero-length write would surface here as a wrong/empty read.
    _, parsed = _call(editor, _bash_body("cat " + path))
    text, is_error = _result(parsed)
    assert is_error is False and text and marker_after in text and marker_before not in text, (
        f"the persisted edit was not readable at {path}: {text!r} "
        f"(want {marker_after!r}, and {marker_before!r} gone)"
    )


def test_i11_presentation_skill_produces_a_valid_pptx_with_text_in_guest():
    """Parity: the guest's presentation-skill toolchain (python-pptx) PRODUCES a
    real, text-bearing .pptx -- a DIFFERENT OOXML type (ppt/slides/slide1.xml)
    and a DIFFERENT skill runtime than the docx (i8) and xlsx (i9) skills. Reads
    the artifact back with zipfile + raw XML only (no python-pptx on the test
    side). This closes the pptx leg of "skills produce their artifact".

    Non-vacuous: assert the deck has a real slide part AND the unique marker text
    is in that slide's XML (not an empty deck or a garbage zip). A guest without
    python-pptx, or a deck missing the marker, reds. Needs a full guest.
    """
    import base64
    import io
    import zipfile

    cid = _cid("i11")
    if not _guest_has_python3(cid):
        pytest.skip("presentation skill needs a full guest (#122) -- SKIP, not a pass.")

    marker = "PPTXMARK" + uuid.uuid4().hex[:8]
    out = "/tmp/i11_deck.pptx"
    gcmd = (
        "python3 - <<PY\n"
        "from pptx import Presentation\n"
        "prs = Presentation()\n"
        "slide = prs.slides.add_slide(prs.slide_layouts[5])\n"
        f"slide.shapes.title.text = {marker!r}\n"
        f"prs.save({out!r})\n"
        "PY\n"
        "base64 -w0 " + out
    )
    _, parsed = _call(cid, _bash_body(gcmd), timeout=150)
    text, is_error = _result(parsed)
    if is_error and text and "pptx" in text and "No module" in text:
        pytest.skip("guest has no python-pptx (presentation skill absent) -- SKIP, not a pass.")
    assert is_error is False and text, f"presentation skill run failed: {text!r}"

    b64 = next((ln for ln in reversed(text.splitlines()) if ln.strip()), "")
    assert len(b64) > 100, f"no base64 artifact streamed from the guest: {text!r}"
    data = base64.b64decode(b64)

    z = zipfile.ZipFile(io.BytesIO(data))
    names = z.namelist()
    assert any("ppt/slides/slide1.xml" in n for n in names), (
        f"the skill artifact is not a valid pptx (no slide part); parts={names}"
    )
    slide_xml = z.read("ppt/slides/slide1.xml").decode(errors="ignore")
    assert marker in slide_xml, (
        f"the produced pptx does not carry the slide text ({marker!r} absent from "
        "the slide) -- the skill did not write the real content"
    )


def test_i12_pdf_skill_produces_a_valid_pdf_with_text_in_guest():
    """Parity: the guest's PDF-skill toolchain (reportlab) PRODUCES a real,
    text-bearing .pdf -- a DIFFERENT artifact format (not OOXML) and a DIFFERENT
    skill runtime than the docx/xlsx/pptx skills. Reads the artifact back as raw
    bytes only (no reportlab on the test side). This closes the pdf leg of
    "skills produce their artifact".

    Non-vacuous: assert the artifact begins with the %PDF header AND the unique
    marker text is present in the PDF bytes. reportlab FlateDecode-compresses the
    content stream by default, which would hide the marker in the raw bytes, so
    the skill run sets pageCompression=0 -- the drawn text then appears verbatim
    in an uncompressed content stream. A guest without reportlab, or a PDF missing
    the marker, reds. Needs a full guest.
    """
    import base64

    cid = _cid("i12")
    if not _guest_has_python3(cid):
        pytest.skip("pdf skill needs a full guest (#122) -- SKIP, not a pass.")

    marker = "PDFMARK" + uuid.uuid4().hex[:8]
    out = "/tmp/i12_doc.pdf"
    gcmd = (
        "python3 - <<PY\n"
        "from reportlab.pdfgen import canvas\n"
        f"c = canvas.Canvas({out!r}, pageCompression=0)\n"
        f"c.drawString(72, 720, {marker!r})\n"
        "c.showPage(); c.save()\n"
        "PY\n"
        "base64 -w0 " + out
    )
    _, parsed = _call(cid, _bash_body(gcmd), timeout=150)
    text, is_error = _result(parsed)
    if is_error and text and "reportlab" in text and "No module" in text:
        pytest.skip("guest has no reportlab (pdf skill absent) -- SKIP, not a pass.")
    assert is_error is False and text, f"pdf skill run failed: {text!r}"

    b64 = next((ln for ln in reversed(text.splitlines()) if ln.strip()), "")
    assert len(b64) > 100, f"no base64 artifact streamed from the guest: {text!r}"
    data = base64.b64decode(b64)

    assert data[:5] == b"%PDF-", (
        f"the skill artifact is not a valid PDF (no %PDF header); head={data[:16]!r}"
    )
    assert marker.encode() in data, (
        f"the produced pdf does not carry the drawn text ({marker!r} absent from "
        "the PDF bytes) -- the skill did not write the real content"
    )
